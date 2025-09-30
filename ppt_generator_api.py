from bson import ObjectId
import certifi
from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
import gridfs
from pydantic import BaseModel
from typing import List
from pymongo import MongoClient
import requests
import json
import demjson3  # pip install demjson3
from pathlib import Path
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from io import BytesIO
import re
from dotenv import load_dotenv
from copy import deepcopy
from groq import Groq
from pptgenerator import build_ppt, get_ppt_from_mongodb, store_ppt_in_mongodb
import os
load_dotenv()

import google.generativeai as genai
gemini_api_key = os.getenv("GEMINI_API_KEY")

genai.configure(api_key=gemini_api_key)

uri = os.getenv("MONGODB_URI")  # example: mongodb+srv://...
client = MongoClient(uri, tlsCAFile=certifi.where())
db = client["ppt_database"]
fs = gridfs.GridFS(db)

# ------------------ Groq Client Setup ------------------ #
# Initialize Groq client with your API key
api_key = os.getenv("GROQ_API_KEY")
client = Groq(
    api_key = api_key
)

# ------------------ Input Schema ------------------ #
class SlideRequest(BaseModel):
    title: str
    slides: int
    model: str  # New field to specify the model
    scrape_from_google: bool = False  # New field to specify if scraping is needed

# ------------------ FastAPI app ------------------ #
origins = [
    "http://localhost:4200",   # Angular dev server
    "http://127.0.0.1:4200",
    "https://pptgenerator-frontend.onrender.com"
    # Add your deployed frontend URL here when hosting Angular
]

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],       # GET, POST, PUT, DELETE...
    allow_headers=["*"],
)

# ------------------ Groq AI Call ------------------ #
def call_groq_ai_system(user_input: str):
    """
    Calls Groq AI endpoint with system/user input and returns JSON slides.
    Replace 'YOUR_GROQ_API_KEY' and endpoint URL with your actual account details.
    """
    chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional presentation writer. Produce a JSON array of slides for the given topic. Return ONLY valid JSON."
                },
                {
                    "role": "user",
                    "content": user_input,
                }
            ],
            # model="llama-3.3-70b-versatile",
            model="meta-llama/llama-4-maverick-17b-128e-instruct",
        )
    
    # Extract the AI-generated content
    try:
        ai_content = chat_completion.choices[0].message.content
    except (AttributeError, IndexError) as e:
        raise HTTPException(status_code=500, detail=f"Invalid Groq AI response structure: {e}")

    # Extract JSON array substring
    match = re.search(r"(\[.*\])", ai_content, re.S)
    if not match:
        raise HTTPException(status_code=500, detail="Could not find JSON array in AI output")

    json_str = match.group(1)

    # --- Sanitize AI output ---
    # 1. Escape unescaped backslashes
    json_str = re.sub(r'\\(?!["\\/bfnrtu])', r'\\\\', json_str)

    # 2. Remove trailing commas before closing brackets/braces
    json_str = re.sub(r',(\s*[\]\}])', r'\1', json_str)

    # 3. Ensure all quotes inside strings are escaped
    # (optional but safer if AI inserts quotes)
    # json_str = re.sub(r'(?<!\\)"', r'\"', json_str)  # use only if needed

    # --- Parse JSON robustly ---
    try:
        # Use demjson3 which can handle non-strict JSON from AI
        slides = demjson3.decode(json_str)
        if not isinstance(slides, list):
            raise HTTPException(status_code=500, detail="AI output is not a JSON array")
        return slides
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI JSON output: {e}")
    

# def call_gemini_ai_system(user_input: str):
#     """
#     Calls Gemini AI endpoint with system/user input and returns JSON slides.
#     Replace 'YOUR_GEMINI_API_KEY' with your actual key.
#     """
#     try:
#         print("Calling Gemini API...")
#         model = genai.GenerativeModel("gemini-2.5-flash")  # or gemini-1.5-pro
#         response = model.generate_content(
#             [
#                 {
#                     "role": "system",
#                     "parts": [
#                         "You are a professional presentation writer. "
#                         "Produce a JSON array of slides for the given topic. "
#                         "Return ONLY valid JSON."
#                     ]
#                 },
#                 {
#                     "role": "user",
#                     "parts": [user_input]
#                 }
#             ]
#         )
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Gemini API call failed: {e}")

#     # Extract AI response
#     try:
#         ai_content = response.text
#     except AttributeError:
#         raise HTTPException(status_code=500, detail="Invalid Gemini response structure")

#     # Extract JSON array substring
#     match = re.search(r"(\[.*\])", ai_content, re.S)
#     if not match:
#         raise HTTPException(status_code=500, detail="Could not find JSON array in AI output")

#     json_str = match.group(1)

#     # --- Sanitize AI output ---
#     json_str = re.sub(r'\\(?!["\\/bfnrtu])', r'\\\\', json_str)   # escape stray backslashes
#     json_str = re.sub(r',(\s*[\]\}])', r'\1', json_str)          # remove trailing commas

#     # --- Parse JSON robustly ---
#     try:
#         slides = demjson3.decode(json_str)
#         if not isinstance(slides, list):
#             raise HTTPException(status_code=500, detail="AI output is not a JSON array")
#         return slides
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Failed to parse AI JSON output: {e}")

def call_gemini_ai_system(user_input: str):
    """
    Calls Gemini AI endpoint with user input and returns JSON slides.
    System prompt is merged into the user role because Gemini only supports 'user' and 'model'.
    """
    try:
        print("Calling Gemini API...")
        model = genai.GenerativeModel("gemini-2.5-flash")  # or gemini-1.5-pro

        system_prompt = (
            "You are a professional presentation writer. "
            "Produce a JSON array of slides for the given topic. "
            "Return ONLY valid JSON."
        )

        response = model.generate_content(
            [
                {
                    "role": "user",
                    "parts": [f"{system_prompt}\n\nTopic: {user_input}"]
                }
            ]
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Gemini API call failed: {e}")

    # Extract AI response
    try:
        ai_content = response.text
    except AttributeError:
        raise HTTPException(status_code=500, detail="Invalid Gemini response structure")

    # Extract JSON array substring
    match = re.search(r"(\[.*\])", ai_content, re.S)
    if not match:
        raise HTTPException(status_code=500, detail="Could not find JSON array in AI output")

    json_str = match.group(1)

    # --- Sanitize AI output ---
    json_str = re.sub(r'\\(?!["\\/bfnrtu])', r'\\\\', json_str)   # escape stray backslashes
    json_str = re.sub(r',(\s*[\]\}])', r'\1', json_str)          # remove trailing commas

    # --- Parse JSON robustly ---
    try:
        slides = demjson3.decode(json_str)
        if not isinstance(slides, list):
            raise HTTPException(status_code=500, detail="AI output is not a JSON array")
        return slides
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI JSON output: {e}")

    

# ------------------ API Endpoint ------------------ #
@app.post("/generate-ppt-slides/")
async def generate_ppt_slides(request: List[SlideRequest]):
    if not request:
        raise HTTPException(status_code=400, detail="No input provided")
    
    # Use first item for simplicity
    topic = request[0].title
    slide_count = request[0].slides
    model = request[0].model
    scrape_from_google = request[0].scrape_from_google


    user_prompt = f"""
    You are a presentation slide generator.

    Topic: { topic }
    Number of slides: { slide_count }

    Instructions:
    - Output ONLY a valid JSON array (no extra text before/after).
    - First array element must only contain the overall presentation title: 
    {{ "title": "Presentation Title" }}

    - Each subsequent slide must be an object with:
    - title (string) → short, clear slide heading
    - content (array) → 4–6 objects, each with:
        - "text": a full, detailed bullet sentence with **keywords in bold** (not just short phrases).
        - "subpoints" (optional array): 4–6 concise sub-bullets expanding on the main point, also with **keywords in bold**.
    - code (object) → if the topic is technical, MUST include:
        - "title": a short label explaining what the code demonstrates
        - "snippet": the snippet should be **multi-line**, detailed, and demonstrate a **practical example** (not trivial). Use `\\n` for line breaks.
    - notes (optional string) → 2–4 sentences for the presenter to elaborate.
    - image_url (optional string) → relevant google search query key string.

    Output format example:

    [
    {{ "title": "Your Presentation Title" }},
    {{
        "title": "Introduction to AI",
        "content": [
        {{
            "text": "**Artificial Intelligence (AI)** is the ability of machines to perform tasks that typically require **human intelligence**.",
            "subpoints": [
            "Focus on **learning algorithms**",
            "Includes **pattern recognition**",
            "Used in **automation** and **decision-making**"
            ]
        }},
        {{
            "text": "AI has transformed **industries** with applications in **healthcare**, **finance**, and **transportation**."
        }}
        ],
        "code": {{
            "title": "Basic Function Example",
            "snippet": "def greet_user(name):\\n    '''This function prints a personalized greeting'''\\n    message = f'Hello, {{name}}! Welcome to Python.'\\n    return message\\n\\nprint(greet_user('Alice'))"
            }},
        "notes": "AI is not a single technology but a field that combines algorithms, data, and computing power.",
        "image_url": "ai-diagram-machine-learning"
    }}
    ]

    Formatting Rules:
    - Each slide must have **exactly 4 or 5 bullet points** in the `content` array.
    - Each bullet should be **detailed** (not just a keyword).
    - Every bullet must contain at least one **bold keyword**.
    - If code is relevant to the topic, include it as a structured object with title + snippet.
    - JSON must be strictly valid.
    """

    if model == "groq":
        slides_json = call_groq_ai_system(user_prompt)
        print("Slides JSON:", slides_json)  # Debugging line
    elif model == "gemini":
        slides_json = call_gemini_ai_system(user_prompt)
        print("Slides JSON:", slides_json)

    with open("debug_slides_json.txt", "w", encoding="utf-8") as f:
        json.dump(slides_json, f, indent=2)

    # call google scrapping for image_url if image_url is present in slide_json
    print(f"Scrape from Google: {scrape_from_google}")
    if scrape_from_google:
        from googlesrapping import scrape_google_images
        for slide in slides_json:
            if "image_url" in slide and slide["image_url"]:
                query = slide["image_url"]
                print(f"Scraping images for query: {query}")
                image_urls = scrape_google_images(query, num_images=5)
                print(f"Scraped {len(image_urls)} image URLs")
                if image_urls:
                    slide["image_url"] = image_urls[0]  # Use the first valid image URL
                    print(f"Found image URL: {slide['image_url']}")
                else:
                    print(f"No valid images found for query: {query}")
                    slide["image_url"] = None  # Clear if no valid image found
            else:
                print("No image_url field in slide or it's empty.")
    

    # Paths
    template_path = "template_iamneo.pptx"
    temp_path = "temp.pptx"
    # split the topic and take first 5 words and join with _
    topic_words = topic.split()[:5]
    topic_short = "_".join(topic_words)
    output_path = f"{topic_short}.pptx"
    print(f"Output path: {output_path}")

    # Build PPT
    # build_ppt(template_path, slides_json, output_path, temp_path)
    # ppt_id = store_ppt_in_mongodb(output_path, Path(output_path).name)
    # get_ppt_from_mongodb(ppt_id, f"downloaded_{Path(output_path).name}")
    # ppt_len = Presentation(output_path)
    # delete temp file and fiel in output path
    Path(temp_path).unlink(missing_ok=True)
    Path(output_path).unlink(missing_ok=True)

    return {"slides": slides_json}


    # return {"message": "PPT generated successfully", "output_file": output_path, "slides_count": len(ppt_len.slides), "ppt_id": str(ppt_id)}

@app.post("/generate-ppt/")
#  request in slide json format
def generate_ppt(request: List[dict]):
    if not request:
        raise HTTPException(status_code=400, detail="No input provided")

    slides_json = request
    topic = slides_json[0].get("title", "Generated_Presentation")
    # Paths
    template_path = "template_iamneo.pptx"
    temp_path = "temp.pptx"
    # output_path = f"{topic.replace(' ', '_')}.pptx"
    topic_words = topic.split()[:5]
    topic_short = "_".join(topic_words)
    output_path = f"{topic_short}.pptx"

    # Build PPT
    build_ppt(template_path, slides_json, output_path, temp_path)
    ppt_id = store_ppt_in_mongodb(output_path, Path(output_path).name)
    # get_ppt_from_mongodb(ppt_id, f"downloaded_{Path(output_path).name}")
    ppt_len = Presentation(output_path)
    # delete temp file and fiel in output path
    Path(temp_path).unlink(missing_ok=True)
    Path(output_path).unlink(missing_ok=True)

    return {"message": "PPT generated successfully", "output_file": output_path, "slides_count": len(ppt_len.slides), "ppt_id": str(ppt_id)}

@app.get("/download/{ppt_id}")
def download_ppt(ppt_id: str):
    try:
        file_id = ObjectId(ppt_id)
        ppt_file = fs.get(file_id)
        print(ppt_file.filename)

        return StreamingResponse(
            ppt_file,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f"attachment; filename={ppt_file.filename}"
            }
        )
    except Exception as e:
        raise HTTPException(status_code=404, detail=f"PPT not found: {e}")