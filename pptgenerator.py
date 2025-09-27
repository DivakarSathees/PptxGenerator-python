from copy import deepcopy
import json
import re
import re
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pymongo import MongoClient
import gridfs
import base64
from dotenv import load_dotenv
from pathlib import Path
import os
import certifi
load_dotenv()

# MongoDB connection
uri = os.getenv("MONGODB_URI")

# client = MongoClient(uri)
client = MongoClient(uri, tlsCAFile=certifi.where())

db = client['ppt_database']       # database name
fs = gridfs.GridFS(db)            # GridFS instance

def store_ppt_in_mongodb(file_path: str, ppt_name: str):
    """
    Stores a PPT file in MongoDB GridFS
    """
    file_path_obj = Path(file_path)
    if not file_path_obj.exists():
        raise FileNotFoundError(f"{file_path} not found")

    # Read binary content
    with open(file_path, "rb") as f:
        file_data = f.read()

    # Store in GridFS
    file_id = fs.put(file_data, filename=ppt_name, contentType="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    print(f"✅ Stored PPT in MongoDB with ID: {file_id}")
    return file_id

def get_ppt_from_mongodb(file_id, save_path):
    data = fs.get(file_id).read()
    with open(save_path, "wb") as f:
        f.write(data)
    print(f"✅ Retrieved PPT from MongoDB: {save_path}")


# ------------------ Helper Functions ------------------ #
def add_bulleted_paragraph(tf, text, level=0):
    """
    Add a paragraph with bullet support and bold/italic formatting.
    level: 0 = main bullet, 1 = sub-bullet
    """
    p = tf.add_paragraph()
    p.level = level
    p.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY
    p.space_after = Pt(5)
    p.bullet = True

    # Temporarily add text for formatting
    p.text = ""
    tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text)
    for token in tokens:
        run = p.add_run()
        if token.startswith("**") and token.endswith("**"):
            run.text = token[2:-2]
            run.font.bold = True
        elif token.startswith("*") and token.endswith("*"):
            run.text = token[1:-1]
            run.font.italic = True
        else:
            run.text = token
        run.font.name = "Calibri"
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0, 0, 0)

    return p

# ------------------ Placeholder Replacement ------------------ #
def replace_placeholders(slide, data):
    """Replace placeholders in a slide."""
    """Replace placeholders in a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # Check if this text frame contains the {content} placeholder
        content_placeholder_found = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == "{content}":
                    content_placeholder_found = True
                    break
            if content_placeholder_found:
                break
        
        if content_placeholder_found and "content" in data:
            # Handle content replacement specially to maintain bullet formatting
            tf = shape.text_frame
            # tf.clear()
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            # Get the original bullet formatting from the first paragraph
            first_para = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            original_level = first_para.level
            
            # Clear all paragraphs except the first one
            for i in range(len(tf.paragraphs) - 1, 0, -1):
                tf.paragraphs[i]._element.getparent().remove(tf.paragraphs[i]._element)
            
            # Clear the first paragraph's text but keep its formatting
            first_para.clear()
            
            # Process content - treat each content item as a separate bullet point
            if data["content"]:
                content_items = []
                
                # Flatten content structure - each main text becomes a bullet point
                for item in data["content"]:
                    if isinstance(item, dict):
                        main_text = item.get("text", "")
                        if main_text:
                            add_bulleted_paragraph(tf, main_text, level=0)

                        for sub in item.get("subpoints", []):
                            add_bulleted_paragraph(tf, sub, level=1)

                    else:
                        add_bulleted_paragraph(tf, str(item), level=0)
                
                # Add first content item to the existing first paragraph
                # if content_items:
                #     first_content = content_items[0]
                #     tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', first_content)
                #     for token in tokens:
                #         run = first_para.add_run()
                #         if token.startswith("**") and token.endswith("**"):
                #             run.text = token[2:-2]
                #             run.font.bold = True
                #         elif token.startswith("*") and token.endswith("*"):
                #             run.text = token[1:-1]
                #             run.font.italic = True
                #         else:
                #             run.text = token
                #         run.font.name = "Calibri"
                #         run.font.size = Pt(22)
                #         run.font.color.rgb = RGBColor(0, 0, 0)
                    
                #     # Set the original level for first paragraph
                #     first_para.level = original_level
                    
                #     # Add remaining content items as new paragraphs with same formatting
                #     for content_text in content_items[1:]:
                #         new_para = tf.add_paragraph()
                #         new_para.level = original_level
                #         new_para.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY
                #         new_para.space_after = Pt(5)
                        
                #         tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', content_text)
                #         for token in tokens:
                #             run = new_para.add_run()
                #             if token.startswith("**") and token.endswith("**"):
                #                 run.text = token[2:-2]
                #                 run.font.bold = True
                #             elif token.startswith("*") and token.endswith("*"):
                #                 run.text = token[1:-1]
                #                 run.font.italic = True
                #             else:
                #                 run.text = token
                #             run.font.name = "Calibri"
                #             run.font.size = Pt(22)
                #             run.font.color.rgb = RGBColor(0, 0, 0)
            
                # Add first content item to the existing first paragraph
                # First content item
                if content_items:
                    first_content = content_items[0]

                    first_para.level = original_level
                    first_para.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY
                    first_para.space_after = Pt(5)
                    first_para.text = ""  # reset

                    # 🔑 Enable bullet for this paragraph
                    first_para.bullet = True

                    tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', first_content)
                    for token in tokens:
                        run = first_para.add_run()
                        if token.startswith("**") and token.endswith("**"):
                            run.text = token[2:-2]
                            run.font.bold = True
                        elif token.startswith("*") and token.endswith("*"):
                            run.text = token[1:-1]
                            run.font.italic = True
                        else:
                            run.text = token
                        run.font.name = "Calibri"
                        run.font.size = Pt(22)
                        run.font.color.rgb = RGBColor(0, 0, 0)

                    # Remaining items
                    for content_text in content_items[1:]:
                        new_para = tf.add_paragraph()
                        new_para.level = original_level
                        new_para.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY
                        new_para.space_after = Pt(5)

                        # 🔑 Enable bullet
                        new_para.bullet = True

                        tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', content_text)
                        for token in tokens:
                            run = new_para.add_run()
                            if token.startswith("**") and token.endswith("**"):
                                run.text = token[2:-2]
                                run.font.bold = True
                            elif token.startswith("*") and token.endswith("*"):
                                run.text = token[1:-1]
                                run.font.italic = True
                            else:
                                run.text = token
                            run.font.name = "Calibri"
                            run.font.size = Pt(22)
                            run.font.color.rgb = RGBColor(0, 0, 0)


            continue 

        # Normal placeholder processing for other placeholders

        # Check if this text frame contains the {content} placeholder
        content_placeholder_found = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == "{content}":
                    content_placeholder_found = True
                    break
            if content_placeholder_found:
                break
        
        if content_placeholder_found and "content" in data:
            # Handle content replacement specially to maintain bullet formatting
            tf = shape.text_frame
            # tf.clear()
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            
            # Get the original bullet formatting from the first paragraph
            first_para = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            original_level = first_para.level
            
            # Clear all paragraphs except the first one
            for i in range(len(tf.paragraphs) - 1, 0, -1):
                tf.paragraphs[i]._element.getparent().remove(tf.paragraphs[i]._element)
            
            # Clear the first paragraph's text but keep its formatting
            first_para.clear()
            
            # Process content - treat each content item as a separate bullet point
            if data["content"]:
                content_items = []
                
                # Flatten content structure - each main text becomes a bullet point
                for item in data["content"]:
                    if isinstance(item, dict):
                        main_text = item.get("text", "")
                        if main_text:
                            add_bulleted_paragraph(tf, main_text, level=0)

                            # content_items.append(main_text)
                        # Add subpoints as separate items with increased indentation
                        for sub in item.get("subpoints", []):
                            add_bulleted_paragraph(tf, sub, level=1)

                        # content_items.extend(subpoints)
                    else:
                        # Handle simple string items
                        add_bulleted_paragraph(tf, str(item), level=0)

                        # content_items.append(str(item))
                
                # Add first content item to the existing first paragraph
                if content_items:
                    first_content = content_items[0]
                    tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', first_content)
                    for token in tokens:
                        run = first_para.add_run()
                        if token.startswith("**") and token.endswith("**"):
                            run.text = token[2:-2]
                            run.font.bold = True
                        elif token.startswith("*") and token.endswith("*"):
                            run.text = token[1:-1]
                            run.font.italic = True
                        else:
                            run.text = token
                        run.font.name = "Calibri"
                        run.font.size = Pt(22)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # Set the original level for first paragraph
                    first_para.level = original_level
                    
                    # Add remaining content items as new paragraphs with same formatting
                    for content_text in content_items[1:]:
                        new_para = tf.add_paragraph()
                        new_para.level = original_level
                        new_para.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY
                        new_para.space_after = Pt(5)
                        
                        tokens = re.split(r'(\*\*.*?\*\*|\*.*?\*)', content_text)
                        for token in tokens:
                            run = new_para.add_run()
                            if token.startswith("**") and token.endswith("**"):
                                run.text = token[2:-2]
                                run.font.bold = True
                            elif token.startswith("*") and token.endswith("*"):
                                run.text = token[1:-1]
                                run.font.italic = True
                            else:
                                run.text = token
                            run.font.name = "Calibri"
                            run.font.size = Pt(22)
                            run.font.color.rgb = RGBColor(0, 0, 0)
            
            continue  # Skip the normal processing for this shape
        
        # Normal placeholder processing for other placeholders
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text.strip()
                print(">>", run.text)
                print(">>", run.text)

                # --- Title ---
                # --- Title ---
                if txt == "{title}" and "title" in data:
                    run.text = data["title"]

                if txt == "codetitle":
                    if "code" in data and isinstance(data["code"], dict) and "title" in data["code"]:
                        run.text = data["code"]["title"]
                        run.font.bold = False
                        run.font.name = "Calibri"
                        run.font.size = Pt(24)
                    else:
                        run.text = ""

                # --- Code ---
                if txt == "{code}":
                    if "code" in data and data["code"]:
                        tf = shape.text_frame
                        tf.clear()
                        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                        p.clear()
                        run = p.add_run()
                        # run.text = data["code"]
                        if isinstance(data["code"], dict):
                            run.text = data["code"].get("snippet", "")
                        else:
                            run.text = data["code"]
                        run.font.name = "Consolas"
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        # fill = shape.fill
                        # fill.solid()
                        # fill.fore_color.rgb = RGBColor(230, 230, 230)
                        p.level = 0
                    else:
                        run.text = ""
                        #  delete the placeholder shape
                        sp = shape.element
                        sp.getparent().remove(sp)

                # --- Notes ---
                if txt == "{notes}" and "notes" in data:
                    if slide.has_notes_slide:
                        slide.notes_slide.notes_text_frame.text = data["notes"]

                from io import BytesIO
                from PIL import Image
                import base64, requests

                if txt == "imageurl":
                    print(data)
                    if "image_url" in data:
                        try:
                            img_url = data["image_url"]

                            if img_url.startswith("data:image"):  # Handle base64-encoded images
                                base64_data = img_url.split(",")[1]
                                image_stream = BytesIO(base64.b64decode(base64_data))
                            else:  # Normal URL, fetch via requests
                                response = requests.get(img_url)
                                print(f"Image fetch status: {response.status_code}")
                                if response.status_code == 200:
                                    image_stream = BytesIO(response.content)
                                else:
                                    run.text = f"status code: {response.status_code} -> {img_url}"
                                    image_stream = None

                            if image_stream:
                                # ✅ Ensure image is in a supported format
                                try:
                                    img = Image.open(image_stream)
                                    if img.format not in ["PNG", "JPEG"]:
                                        converted_stream = BytesIO()
                                        img.convert("RGB").save(converted_stream, format="PNG")
                                        converted_stream.seek(0)
                                        image_stream = converted_stream
                                except Exception as e:
                                    print(f"⚠️ Error converting image: {e}")

                                run.text = ""
                                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

                                # remove original placeholder
                                sp = shape.element
                                sp.getparent().remove(sp)

                        except Exception as e:
                            print(f"⚠️ Could not add image: {e}")
                    else:
                        run.text = ""


                # if txt == "imageurl":
                #     print(data)
                #     if "image_url" in data:
                #         try:
                #             img_url = data["image_url"]

                #             if img_url.startswith("data:image"):  # Handle base64-encoded images
                #                 # Extract only the base64 part after the comma
                #                 base64_data = img_url.split(",")[1]
                #                 image_stream = BytesIO(base64.b64decode(base64_data))
                #             else:  # Normal URL, fetch via requests
                #                 response = requests.get(img_url)
                #                 print(f"Image fetch status: {response.status_code}")
                #                 if response.status_code == 200:
                #                     image_stream = BytesIO(response.content)
                #                 else:
                #                     run.text = f"status code: {response.status_code} -> {img_url}"
                #                     image_stream = None

                #             if image_stream:
                #                 run.text = ""
                #                 left, top, width, height = shape.left, shape.top, shape.width, shape.height
                #                 slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

                #                 # remove original placeholder
                #                 sp = shape.element
                #                 sp.getparent().remove(sp)

                #         except Exception as e:
                #             print(f"⚠️ Could not add image: {e}")
                #     else:
                #         run.text = ""

# ------------------ Slide Duplication ------------------ #
# ------------------ Slide Duplication ------------------ #
# def duplicate_slide(prs, slide):
#     """Duplicate a slide while excluding placeholders."""
#     """Duplicate a slide while excluding placeholders."""
#     slide_layout = slide.slide_layout
#     new_slide = prs.slides.add_slide(slide_layout)

#     # remove placeholders
#     # remove placeholders
#     for shape in list(new_slide.shapes):
#         if shape.is_placeholder:
#             sp = shape.element
#             sp.getparent().remove(sp)

#     # copy original shapes
#     # copy original shapes
#     for shape in slide.shapes:
#         new_el = deepcopy(shape.element)
#         new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

#     return new_slide

# def build_ppt(template_path, slides_json, output_path, temp_path):
#     prs1 = Presentation(template_path)

#     # Step 1: Define layouts (assuming index 1 = content, 2 = code)
#     content_layout_index = 1   # 2nd slide in template
#     code_layout_index = 2      # 3rd slide in template

#     # Step 2: Build expanded slide plan
#     expanded_slides = []
#     for slide_data in slides_json:
#         if "code" in slide_data and slide_data["code"]:
#             expanded_slides.append({"layout": content_layout_index, "data": slide_data, "mode": "content"})
#             expanded_slides.append({"layout": code_layout_index, "data": slide_data, "mode": "code"})
#         else:
#             expanded_slides.append({"layout": content_layout_index, "data": slide_data, "mode": "content"})

#     # Step 3: Ensure enough slides exist by duplicating the right layout
#     template_slide_count = len(prs1.slides)
#     for idx in range(len(expanded_slides)):
#         if idx >= template_slide_count:
#             layout_index = expanded_slides[idx]["layout"]
#             duplicate_slide(prs1, prs1.slides[layout_index])

#     prs1.save(temp_path)
#     prs = Presentation(temp_path)

#     # Step 4: Fill slides
#     for idx, slide_info in enumerate(expanded_slides):
#         slide = prs.slides[idx]
#         if slide_info["mode"] == "code":
#             code_data = {
#                 "title": "Example: " + slide_info["data"]["title"],
#                 "content": [],
#                 "code": slide_info["data"]["code"],
#                 "notes": slide_info["data"].get("notes", "")
#             }
#             replace_placeholders(slide, code_data)
#         else:
#             content_data = dict(slide_info["data"])
#             content_data["code"] = ""   # 🚫 clear code for non-code slides
#             print(">>", content_data)
#             replace_placeholders(slide, content_data)

#     prs.save(output_path)
#     print(f"✅ Final PPT created: {output_path}")

def split_code_into_chunks(code_str, max_lines=25):
    """Split a code snippet into chunks of max_lines each."""
    lines = code_str.splitlines()
    chunks = []
    for i in range(0, len(lines), max_lines):
        chunk = "\n".join(lines[i:i+max_lines])
        chunks.append(chunk)
    return chunks

def duplicate_slide(prs, slide):
    """Duplicate a slide while excluding placeholders."""
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # remove placeholders
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # copy original shapes
    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide

def build_ppt(template_path, slides_json, output_path, temp_path):
    prs1 = Presentation(template_path)

    # Step 1: Define layouts (assuming index 1 = content, 2 = code)
    content_layout_index = 1   # 2nd slide in template
    code_layout_index = 2      # 3rd slide in template

    # Step 2: Build expanded slide plan
    expanded_slides = []
    for slide_data in slides_json:
        if "code" in slide_data and slide_data["code"]:
            # Split code into chunks of 25 lines
            code_chunks = split_code_into_chunks(slide_data["code"]["snippet"], max_lines=25)

            # First: normal content slide
            expanded_slides.append({"layout": content_layout_index, "data": slide_data, "mode": "content"})

            # Then: one slide per code chunk
            for idx, chunk in enumerate(code_chunks):
                chunk_data = dict(slide_data)
                chunk_data["code"] = {
                    "title": slide_data["code"]["title"] + (f" (Part {idx+1})" if len(code_chunks) > 1 else ""),
                    "snippet": chunk
                }
                expanded_slides.append({"layout": code_layout_index, "data": chunk_data, "mode": "code"})
        else:
            expanded_slides.append({"layout": content_layout_index, "data": slide_data, "mode": "content"})

    # Step 3: Ensure enough slides exist by duplicating the right layout
    template_slide_count = len(prs1.slides)
    for idx in range(len(expanded_slides)):
        if idx >= template_slide_count:
            layout_index = expanded_slides[idx]["layout"]
            duplicate_slide(prs1, prs1.slides[layout_index])

    prs1.save(temp_path)
    prs = Presentation(temp_path)

    # Step 4: Fill slides
    for idx, slide_info in enumerate(expanded_slides):
        slide = prs.slides[idx]
        if slide_info["mode"] == "code":
            code_data = {
                "title": "Example: " + slide_info["data"]["title"],
                "content": [],
                "code": slide_info["data"]["code"],
                "notes": slide_info["data"].get("notes", "")
            }
            replace_placeholders(slide, code_data)
        else:
            content_data = dict(slide_info["data"])
            content_data["code"] = ""   # 🚫 clear code for non-code slides
            replace_placeholders(slide, content_data)

    prs.save(output_path)
    print(f"✅ Final PPT created: {output_path}")



# ------------------ Main ------------------ #
if __name__ == "__main__":
    build_ppt("template_iamneo.pptx", "slides.json", "Cloud_Trends_2025.pptx", "temp.pptx")